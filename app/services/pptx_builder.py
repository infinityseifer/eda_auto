"""PowerPoint builder for Auto-EDA MVP.

Creates a presentation with:
- Title + Executive summary
- Data Overview
- Key Drivers (bullets)
- Anomalies & Caveats (bullets)
- Recommendations (bullets)
- Chart gallery (up to N images)
"""
# (full minimal theming version)
# app/services/pptx_builder.py
from __future__ import annotations

from pathlib import Path
from typing import Dict

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ✅ use centralized settings to avoid path drift
from app.core.config import settings


def _hex_to_rgb(hexstr: str) -> tuple[int, int, int]:
    s = hexstr.strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    s = s[:6].ljust(6, "0")
    return tuple(int(s[i : i + 2], 16) for i in (0, 2, 4))

def _palette(theme: str, accent_hex: str) -> Dict[str, tuple[int, int, int]]:
    light = {"bg": (255, 255, 255), "text": (17, 17, 17)}
    dark  = {"bg": (17, 17, 17),  "text": (255, 255, 255)}
    base = light if str(theme).lower() != "dark" else dark
    acc = _hex_to_rgb(accent_hex or "#1f77b4")
    return {"bg": base["bg"], "text": base["text"], "accent": acc}

def _fill_shape_rgb(shape, rgb: tuple[int, int, int]):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

def _set_text(run, text: str, size: int, rgb: tuple[int, int, int], bold=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)

def _add_title_slide(prs: Presentation, title: str, subtitle: str, pal) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    _fill_shape_rgb(bg, pal["bg"])
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(9), Inches(1.2))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p.runs[0] if p.runs else p.add_run(), title, 44, pal["text"], bold=True)
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(3.0), Inches(2.5), Inches(0.12))
    _fill_shape_rgb(line, pal["accent"])
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(9), Inches(0.8))
    p2 = tx2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    _set_text(p2.runs[0] if p2.runs else p2.add_run(), subtitle, 20, pal["text"])

def _add_section_slide(prs: Presentation, heading: str, body_lines: list[str], pal) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    _fill_shape_rgb(bg, pal["bg"])
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(9), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _set_text(p.runs[0] if p.runs else p.add_run(), heading, 28, pal["text"], bold=True)
    bar = slide.shapes.add_shape(1, Inches(0.8), Inches(1.75), Inches(1.8), Inches(0.08))
    _fill_shape_rgb(bar, pal["accent"])
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(9.2), Inches(4.5))
    tf = tx2.text_frame
    tf.clear()
    for i, line in enumerate(body_lines or []):
        par = tf.add_paragraph() if i else tf.paragraphs[0]
        par.level = 0
        par.alignment = PP_ALIGN.LEFT
        run = par.add_run()
        _set_text(run, str(line), 18, pal["text"])

def _add_chart_slide(prs: Presentation, title: str, img_path: str, pal) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    _fill_shape_rgb(bg, pal["bg"])
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(9), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    _set_text(p.runs[0] if p.runs else p.add_run(), title, 22, pal["text"], bold=True)
    left, top, w = Inches(0.8), Inches(1.4), Inches(8.8)
    img = Path(img_path)
    if pal["bg"] == (17, 17, 17):  # dark theme
        mat = slide.shapes.add_shape(1, left, top, w, Inches(5.0))
        _fill_shape_rgb(mat, (255, 255, 255))
        slide.shapes.add_picture(str(img), left + Inches(0.15), top + Inches(0.15), width=w - Inches(0.3))
    else:
        slide.shapes.add_picture(str(img), left, top, width=w)

def build_presentation(eda: dict, narrative: dict, out_dir: str | None = None, *, theme: str = "light", color: str = "#1f77b4") -> str:
    """
    Build the PPTX and save it under settings.STORAGE_DIR/reports.
    Returns the filename (not the absolute path) so API download stays consistent.
    """
    pal = _palette(theme, color)
    prs = Presentation()

    # Title + narrative slides
    _add_title_slide(prs, "Auto EDA & Storytelling", f"Dataset: {eda.get('dataset_id','')}", pal)
    _add_section_slide(prs, "Executive Summary", [narrative.get("executive_summary", "")], pal)
    _add_section_slide(prs, "Data Overview", (narrative.get("data_overview", "") or "").splitlines(), pal)
    _add_section_slide(prs, "Key Drivers", narrative.get("key_drivers", []), pal)
    _add_section_slide(prs, "Anomalies & Caveats", narrative.get("anomalies", []), pal)
    _add_section_slide(prs, "Recommendations", narrative.get("recommendations", []), pal)

    # Chart slides
    for title, img in (eda.get("charts") or {}).items():
        _add_chart_slide(prs, title, img, pal)

    # ✅ Canonical save location
    reports_dir = Path(settings.STORAGE_DIR) / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    name = f"report_{eda.get('dataset_id','dataset')}_{theme}.pptx"
    out_path = reports_dir / name
    prs.save(out_path)

    # Return filename only
    return name

# Optional back-compat alias
def build_pptx(eda, narrative, out_dir: str | None = None, *, theme: str = "light", color: str = "#1f77b4") -> str:
    return build_presentation(eda=eda, narrative=narrative, out_dir=out_dir, theme=theme, color=color)
